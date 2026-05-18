# -*- coding: utf-8 -*-
import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / ".tools" / "python-pptx"))

from pptx import Presentation


SOURCE = Path(r"C:\Users\cndqj\Downloads\Pawsitive____community_light.pptx")
OUTPUT = Path(r"C:\Users\cndqj\Downloads\Pawsitive____community_flow.pptx")
BACKUP = Path(r"C:\Users\cndqj\Downloads\Pawsitive____community_light_backup_before_flow.pptx")

REPLACEMENTS = {
    "품종 탐색부터 결제, 산책, 건강 분석, 상담, 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.":
    "품종 탐색부터 결제, 산책, 기록 공유, 건강 분석, 상담, 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.",
    "경로 · 시간 · 거리 · 공유 기록":
    "경로 · 시간 · 거리 · 커뮤니티 공유",
    "거리·시간·경로·공유":
    "거리·시간·경로·공유 게시",
    "구현 영상 04 · 산책 → 건강 → 상담":
    "구현 영상 04 · 산책 → 공유 → 건강 → 상담",
    "산책 기록이 건강 분석과 상담의 입력값으로 이어집니다.":
    "산책 기록은 공유되고, 건강 분석과 상담의 입력값으로 이어집니다.",
    "완료된 산책 기록":
    "완료된 산책 기록 공유",
    "거리, 시간, 경로를 기록하고 필요하면 커뮤니티에 가볍게 공유할 수 있습니다.":
    "산책 카드로 저장된 거리·시간·경로를 커뮤니티 게시글에 첨부해 공유합니다.",
    "영상 04 · 산책 기록 → 건강 분석 → AI 상담":
    "영상 04 · 산책 기록 → 공유 → 건강 분석 → AI 상담",
    "DEMO 04 · /WALK · /HEALTH · /AI":
    "DEMO 04 · /WALK · /COMMUNITY · /HEALTH · /AI",
    "산책 데이터가 공유와 상담 맥락으로 흐릅니다":
    "산책 기록 → 커뮤니티 공유 → 건강 분석 → AI 상담",
    "경로 · 시간 · 거리 · 공유":
    "경로 · 시간 · 거리 · 커뮤니티 공유",
    "보호자는 같은 설명을 반복하지 않고, 전문가는 상담 전 반려견의 맥락을 먼저 이해할 수 있습니다.":
    "보호자는 산책 루틴을 서로 참고하고, 전문가는 상담 전 반려견 맥락을 빠르게 이해할 수 있습니다.",
    "산책 기록과 건강 리포트를 상담 전 공유해 진료의 시작 시간을 단축합니다.":
    "산책 기록과 건강 리포트를 상담 전 전달해 진료·훈련 상담의 시작 시간을 줄입니다.",
}

TEXT_OVERRIDES = {
    (9, 7): "산책 기록은 커뮤니티에 공유되고, 건강 분석과 상담의 입력값으로 이어집니다.",
}


def replace_text_in_shape(shape, counts):
    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            for old, new in REPLACEMENTS.items():
                if old in text:
                    text = text.replace(old, new)
                    counts[old] += 1
            run.text = text


def overwrite_shape_text(shape, text):
    if not getattr(shape, "has_text_frame", False):
        return
    wrote = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not wrote:
                run.text = text
                wrote = True
            else:
                run.text = ""
    if not wrote:
        shape.text_frame.text = text


def main():
    if not BACKUP.exists():
        shutil.copy2(SOURCE, BACKUP)

    deck = Presentation(str(SOURCE))
    counts = {old: 0 for old in REPLACEMENTS}

    for slide in deck.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, counts)

    for (slide_no, shape_idx), text in TEXT_OVERRIDES.items():
        overwrite_shape_text(deck.slides[slide_no - 1].shapes[shape_idx], text)

    deck.save(str(OUTPUT))

    print(str(OUTPUT))
    for old, count in counts.items():
        print(f"{count}\t{old}")


if __name__ == "__main__":
    main()
