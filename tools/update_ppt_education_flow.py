import os
import shutil
import sys
from pathlib import Path

sys.path.insert(0, os.path.abspath(os.path.join(".tools", "python-pptx")))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


SRC = Path.home() / "Downloads" / "Pawsitive____community_flow.pptx"
OUT = Path.home() / "Downloads" / "Pawsitive____community_flow_education.pptx"


REPLACEMENTS = {
    "품종 탐색부터 결제, 산책, 기록 공유, 건강 분석, 상담, 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.":
        "품종 탐색부터 결제, 산책, 기록 공유, 건강 분석, AI 상담, 교육센터 학습, 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.",
    "전문가 매칭":
        "교육센터 · 전문가 매칭",
    "병원·훈련사 · 맥락 인계":
        "학습 콘텐츠 · 병원/훈련사 인계",
    "구현 영상 04 · 산책 → 공유 → 건강 → 상담":
        "구현 영상 04 · 산책 → 공유 → 건강 → 상담 → 학습",
    "산책 기록은 커뮤니티에 공유되고, 건강 분석과 상담의 입력값으로 이어집니다.":
        "산책 기록은 공유되고, 건강 분석·AI 상담·교육센터로 이어집니다.",
    "RAG 상담":
        "RAG 상담 · 교육센터 연계",
    "\"최근 산책량이 줄었는데 괜찮을까요?\" 같은 질문에 내부 지식 + 산책·건강 맥락을 합쳐 답변합니다.":
        "\"최근 산책량이 줄었는데 괜찮을까요?\" 같은 질문에 내부 지식 + 산책·건강 맥락을 합쳐 답변하고, 필요한 관리 지식은 교육센터에서 이어서 학습합니다.",
    "DEMO 04 · /WALK · /COMMUNITY · /HEALTH · /AI":
        "DEMO 04 · /WALK · /COMMUNITY · /HEALTH · /AI · /EDUCATION",
    "영상 04 · 산책 기록 → 공유 → 건강 분석 → AI 상담":
        "영상 04 · 공유 → 건강 → 상담 → 교육",
    "산책 기록 → 커뮤니티 공유 → 건강 분석 → AI 상담":
        "산책 기록 공유 → 건강 분석 → AI 상담 → 교육센터 학습",
    "RAG 상담 · 교육센터 연계 (Hybrid 70/30)":
        "RAG상담(Hybrid70/30)",
    "AI 상담 인터페이스":
        "AI상담/교육",
    "건강 문서":
        "교육DB 44개",
    "전문가 상담 기록":
        "교육 · 전문가 상담 기록",
    "병원 · 훈련사 인계":
        "교육 진도 · 병원/훈련사 인계",
    "AI 상담 이력":
        "AI 상담 · 학습 이력",
    "맥락 · 질문 · 답변":
        "질문 · 답변 · 학습 연결",
    "보호자는 산책 루틴을 서로 참고하고, 전문가는 상담 전 반려견 맥락을 빠르게 이해할 수 있습니다.":
        "보호자는 산책 루틴과 교육 콘텐츠를 함께 참고하고, 전문가는 상담 전 반려견 맥락을 빠르게 이해할 수 있습니다.",
    "산책 기록과 건강 리포트를 상담 전 전달해 진료·훈련 상담의 시작 시간을 줄입니다.":
        "산책 기록, 건강 리포트, 교육 이력을 상담 전 전달해 진료·훈련 상담의 시작 시간을 줄입니다.",
    "RAG, GPS, 실시간 매칭, 건강 분석을 4-Layer로 통합했습니다.":
        "RAG, GPS, 실시간 매칭, 건강 분석, 교육 콘텐츠를 4-Layer로 통합했습니다.",
    "품종 추천부터 상담까지 end-to-end 흐름이 실제 동작합니다.":
        "품종 추천부터 산책, 건강 분석, AI 상담, 교육센터까지 end-to-end 흐름이 실제 동작합니다.",
    "품종 탐색부터 결제, 산책, 기록 공유, 건강 분석, AI 상담, 교육센터 학습, 교육센터 · 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.":
        "품종 탐색부터 결제, 산책, 기록 공유, 건강 분석, AI 상담, 교육센터 학습, 전문가 매칭까지. Pawsitive에서 산책은 사라지는 이벤트가 아니라 다음 케어 판단을 위한 데이터가 됩니다.",
}


def replace_text(prs: Presentation) -> int:
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text
                    changed = original
                    for old, new in REPLACEMENTS.items():
                        changed = changed.replace(old, new)
                    if changed != original:
                        run.text = changed
                        count += 1
    return count


def add_demo04_note(prs: Presentation) -> None:
    slide = prs.slides[8]
    note = slide.shapes.add_textbox(Inches(0.78), Inches(6.22), Inches(4.85), Inches(0.36))
    text_frame = note.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "추가 컷 추천: AI 상담 답변 후 교육센터에서 훈련/건강 콘텐츠를 확인하고 퀴즈·진도율이 올라가는 장면"
    p.font.size = Pt(7.5)
    p.font.name = "Pretendard"
    p.font.color.rgb = RGBColor(90, 82, 72)


def main() -> None:
    shutil.copy2(SRC, OUT)
    prs = Presentation(OUT)
    count = replace_text(prs)
    prs.save(OUT)
    print(OUT)
    print(f"updated={count}")


if __name__ == "__main__":
    main()
